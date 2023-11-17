import base64
import datetime
import glob
import io
import json
import os
import random
import uuid
import copy

from fastapi import APIRouter, Depends, HTTPException, Response
from fastapi.responses import FileResponse, StreamingResponse
from py2neo import Node, Transaction, NodeMatcher, Relationship as Neo4jRelationship

from app.api.knowledge.models import *
from app.api.knowledge.utils import *
from app.api.nlp.nlp import add_text, clear_cache_texts
from app.graph.connector2 import neodbsession
from app.api.ppt.utils import chrome
import hashlib
from app.api.ppt.connector import minio_client
from PIL import Image
from pptx import Presentation
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.util import Inches
from pptx.shapes.autoshape import Shape
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
import math

router = APIRouter()
url = "/ppt"


def dict_hash(dictionary: Dict) -> str:
    """MD5 hash of a dictionary."""
    dhash = hashlib.md5()
    # We need to sort arguments so {'a': 1, 'b': 2} is
    # the same as {'b': 2, 'a': 1}
    encoded = json.dumps(dictionary, sort_keys=True).replace(" ", "").encode()
    dhash.update(encoded)
    ret = dhash.hexdigest()
    return ret


def process_request(
    url, width, height, wait, only_hash=True, force=False, return_thumb=True
):
    global chrome

    hash_code = dict_hash(dict(url=url, width=width, height=height))
    response = None
    pic = None
    if not force:
        try:
            suffix = "_thumb.png" if return_thumb else ".png"
            response = minio_client.get_object("backup", hash_code + suffix)
            if only_hash:
                return hash_code
            return response.read()
        except Exception as e:
            print("ERROR", e)
        finally:
            if response:
                response.close()
                response.release_conn()

    pic = chrome.take_pic(url, width, height, wait)
    pic_bytes = io.BytesIO(pic)

    result = minio_client.put_object(
        "backup",
        hash_code + ".png",
        pic_bytes,
        length=-1,
        part_size=10 * 1024 * 1024,
        content_type="image/png",
    )

    thumb = Image.open(pic_bytes)
    thumb.thumbnail((300, 300))

    in_mem_thumb = io.BytesIO()
    thumb.save(in_mem_thumb, format="png")
    in_mem_thumb.seek(0)

    result = minio_client.put_object(
        "backup",
        hash_code + "_thumb.png",
        in_mem_thumb,
        length=-1,
        part_size=10 * 1024 * 1024,
        content_type="image/png",
    )

    print("uploaded backup", result)

    if only_hash:
        return hash_code

    if return_thumb:
        in_mem_thumb.seek(0)
        return in_mem_thumb.read()
    else:
        pic_bytes.seek(0)
        return pic_bytes.read()


@router.get(url, tags=["ppt"])
def example():
    pic = process_request(
        "https://qol.vav.aknakos.com/?tree=0b1b6a0b1a0b0b0&map=LIFESAT&maximized=0&split=1&bo=linkert&oo=linkert&mo=norm_mean&weighted=0&mapAreaType=postal&lat=43.297198404646366&long=-61.309526506442204&zoom=7",
        2495,
        1213,
        3,
        only_hash=False,
    )

    return Response(content=pic, media_type="image/png")


class PictureRequest(BaseModel):
    href: str
    width: int
    height: int
    wait: int = 4
    force: bool = False


@router.post(url, responses={200: {"content": {"image/png": {}}}}, tags=["ppt"])
def take_pic(request: PictureRequest):
    pic = process_request(
        request.href,
        request.width,
        request.height,
        request.wait,
        only_hash=False,
        force=request.force,
    )

    return Response(content=pic, media_type="image/png")


class PptRequestPicture(BaseModel):
    href: str
    width: int
    height: int
    wait: int
    label: str
    type: str
    shapes: List[Dict]
    data: Dict


class PptRequest(BaseModel):
    ppts: List[PptRequestPicture]


@router.post(
    url + "/generate",
    responses={
        200: {
            "content": {
                "application/vnd.openxmlformats-officedocument.presentationml.presentation": {}
            }
        }
    },
    tags=["ppt"],
)
def generate_ppt(request: PptRequest):
    print(request)
    pics = [
        process_request(
            x.href,
            x.width,
            x.height,
            x.wait,
            only_hash=False,
            force=False,
            return_thumb=False,
        )
        for x in request.ppts
    ]

    prs = Presentation()

    for i, req in enumerate(request.ppts):
        text = req.data["text"]
        # ret = self.driver.get_screenshot_as_base64()
        pic = io.BytesIO(pics[i])
        last = i == len(request.ppts) - 1
        first = i == 0
        intention = req.label == "H_SEQUENCE" and req.type == "intention"
        insight = req.label == "H_SEQUENCE" and req.type == "insight"
        interaction = req.label == "M_SEQUENCE"

        picture_slide_layout = prs.slide_layouts[8]
        slide = prs.slides.add_slide(picture_slide_layout)
        picture_placeholder = slide.placeholders[1]
        shapes = slide.shapes
        title = slide.placeholders[0]
        s_text = slide.placeholders[2]
        title.text = ""

        if req.label == "H_SEQUENCE":
            if last:
                title.text += "Final "
            elif first:
                title.text += "Initial "
            if insight:
                title.text += "Insight: "
            else:
                title.text += "Intention: "

            # title.text += text
            s_text.text += text
        elif interaction:
            title.text += "Interaction: "
            # title.text += text
            s_text.text += text
        else:
            title.text = text

        title.left, s_text.left = Inches(1), Inches(1)
        title.width, s_text.width = prs.slide_width - Inches(
            2
        ), prs.slide_width - Inches(2)
        title.height, s_text.height = Inches(0.5), Inches(1)
        title.top, s_text.top = prs.slide_height - Inches(3), prs.slide_height - Inches(
            2
        )

        # slide.placeholders[0].text = 'Final Insight: ' + p['text'] if p.has_label('H_UPDATE') else p['event']
        # slide.placeholders[0].text = text
        # slide.placeholders[1].text = title.text
        # slide.placeholders[2].text = title.text

        req.height = req.height * 1.2
        deltax = prs.slide_width / int(req.width)
        deltay = deltax
        picture_placeholder.left = 0
        picture_placeholder.top = 0
        picture_placeholder.width = prs.slide_width
        picture_placeholder.height = int(req.height * deltay)
        picture_placeholder.crop_top = 0
        picture_placeholder.crop_left = 0
        picture_placeholder.crop_bottom = 0
        picture_placeholder.crop_right = 0

        pic_shape = picture_placeholder.insert_picture(pic)
        pic_shape.left = 0
        pic_shape.top = 0
        pic_shape.width = prs.slide_width
        pic_shape.height = int(req.height * deltay)
        pic_shape.crop_top = 0
        pic_shape.crop_left = 0
        pic_shape.crop_bottom = 0
        pic_shape.crop_right = 0

        title.top, s_text.top = int(req.height * deltay) + Inches(0.3), int(
            req.height * deltay
        ) + Inches(0.85)
        print(deltay)
        for s in req.shapes:
            if s["type"] == "circle":
                shape = shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.OVAL,
                    (s["x"] - s["rx"]) * deltax,
                    (s["y"] - s["ry"]) * deltay,
                    (s["rx"] * 2) * deltax,
                    (s["ry"] * 2) * deltay,
                )
            else:
                x1, x2, y1, y2 = (
                    s["x1"] * deltax,
                    s["x2"] * deltax,
                    (s["y1"]) * deltay,
                    (s["y2"]) * deltay,
                )
                # if not takepictures:
                #     shape2 = shapes.add_shape(
                #         MSO_AUTO_SHAPE_TYPE.OVAL,
                #         x1 - Inches(1) / 2,
                #         y1 - Inches(1) / 2,
                #         Inches(1),
                #         Inches(1),
                #     )
                #     shape1 = shapes.add_shape(
                #         MSO_AUTO_SHAPE_TYPE.OVAL,
                #         x2 - Inches(1) / 2,
                #         y2 - Inches(1) / 2,
                #         Inches(1),
                #         Inches(1),
                #     )
                #     shape1.fill.background()
                #     shape2.fill.background()

                w = Inches(0.2)
                top = y1
                left = x1
                dist = math.hypot(x2 - x1, y2 - y1) + (deltax * 10)
                ang = -(math.atan2(y2 - y1, x2 - x1))
                deltatop = math.sin(ang) * (dist - w * math.cos(ang)) / 2
                deltaleft = math.tan((math.pi * 2 - ang) / 2) * (
                    deltatop - w * math.sin(ang) / 2
                )
                print(top, deltatop, math.degrees(ang))
                shape = shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
                    left + deltaleft,
                    top - deltatop,
                    dist,
                    w,
                )
                shape.rotation = -math.degrees(ang)
            shape.fill.background()

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    return StreamingResponse(
        content=output,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

    # return Response(
    #     content=output.read(),
    #     media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    # )

    # global chrome
    # pic = chrome.take_pic(request.href, request.width, request.height, request.wait)

    # return Response(content=pic, media_type="image/png")
